from __future__ import annotations

import asyncio
import logging
import shutil
from pathlib import Path
from typing import Callable

from . import jobs as jobs_mod
from .combine import combine as do_combine
from .lang.detect import detect_language
from .lang.rtl import is_rtl
from .llm.client import LLMClient
from .rtl import apply_rtl
from .schemas import DocFormat, OutputMode, TranslationJob
from .translators.registry import get_translator
from .utils.errors import TranslationError
from .utils.io import job_workspace

log = logging.getLogger(__name__)


async def translate_document(
    job: TranslationJob,
    progress: Callable[[float, str], None] | None = None,
) -> tuple[Path, DocFormat]:
    raise RuntimeError("Use run_job() which manages workspace lifetime.")


async def run_job(
    job: TranslationJob,
    progress: Callable[[float, str], None] | None = None,
    job_id: str | None = None,
) -> Path:
    """Execute the job end-to-end and return a path to the deliverable.

    The returned path is OUTSIDE the temp workspace (moved to a stable location)
    so callers can stream it to the user without racing cleanup.
    """
    from .config import settings

    def _p(frac: float, msg: str) -> None:
        if progress:
            progress(frac, msg)

    src = job.src_path
    fmt = DocFormat.from_path(src)
    translator = get_translator(src)

    with job_workspace() as wd:
        if fmt is DocFormat.PDF:
            from .translators.pdf_ocr import maybe_ocr_pdf

            _p(0.01, "Running OCR")
            src = maybe_ocr_pdf(src, wd)

        # 1. Detect source language (optional, purely for prompting context).
        _p(0.02, "Detecting source language")
        source_lang = job.source_lang
        if not source_lang:
            # Use a quick text snippet from the doc for detection.
            snippet = _sample_text(src, fmt)
            if snippet:
                source_lang = detect_language(snippet)
        log.info("source_lang=%s target_lang=%s", source_lang, job.target_lang)

        # 2. Extract segments.
        _p(0.08, "Extracting text segments")
        segments = translator.extract(src)
        log.info("extracted %d segments", len(segments))

        # 3. Call LLM to translate.
        client = LLMClient(
            model=job.model_string,
            api_key=job.api_key,
            api_base=job.api_base,
            temperature=job.temperature,
        )

        # 3a. Build a document-wide context brief BEFORE translating, so every
        # chunk is translated with the same domain vocabulary. Without this,
        # domain-specific terms like "sprint" (Scrum) or "pipeline" (DevOps)
        # risk being translated literally and out of context. One summary call
        # per job; cost stays negligible even on very long docs because we
        # sample text from start/middle/end rather than sending everything.
        _p(0.10, "Reading document for context")
        context = await client.summarize_document(
            segments,
            target_lang=job.target_lang,
            source_lang=source_lang,
        )
        if context:
            log.info("using document context (%d chars)", len(context))

        def _llm_progress(f: float, msg: str) -> None:
            _p(0.12 + f * 0.66, msg)

        await client.translate_segments(
            segments,
            target_lang=job.target_lang,
            source_lang=source_lang,
            max_chunk_tokens=job.max_chunk_tokens,
            progress=_llm_progress,
            context=context,
        )

        # 4. Reinsert into format-preserving output.
        _p(0.80, "Rebuilding document")
        translated_path = wd / f"translated{src.suffix}"
        if isinstance(translator, _import_pdf_cls()):
            # PdfTranslator.reinsert accepts target_lang for inline RTL handling.
            translator.reinsert(src, segments, translated_path, target_lang=job.target_lang)
        else:
            translator.reinsert(src, segments, translated_path)

        # 5. Apply RTL (non-PDF formats).
        if is_rtl(job.target_lang):
            _p(0.88, "Applying RTL direction")
            apply_rtl(translated_path, fmt)

        # Cache the pre-combine ("raw") artifact so future requests for a
        # different output_mode can combine it without re-running the LLM.
        if job_id:
            raw_name = f"__raw{src.suffix}"
            raw_stable = jobs_mod.output_path(job_id, raw_name)
            raw_stable.unlink(missing_ok=True)
            _link_or_copy(translated_path, raw_stable)
            jobs_mod.update_status(job_id, raw_name=raw_name)

        # 6. Combine if needed.
        mode = job.output_mode
        if mode is OutputMode.ORIGINAL:
            final = src
            ext = src.suffix
        elif mode is OutputMode.TRANSLATED:
            final = translated_path
            ext = src.suffix
        else:
            _p(0.93, f"Combining ({mode.value})")
            combined_suffix = src.suffix
            combined_path = wd / f"combined{combined_suffix}"
            final = do_combine(
                src_path=src,
                translated_path=translated_path,
                out_path=combined_path,
                fmt=fmt,
                mode=mode,
                rtl=is_rtl(job.target_lang),
            )
            ext = final.suffix  # may become .pdf for DOCX/PPTX horizontal

        # 7. Move deliverable to a stable location. If a job_id is provided,
        # persist to the job directory (so it survives workspace cleanup and
        # becomes retrievable via GET /api/jobs/{id}/download). Otherwise use
        # the generic out dir for backwards-compat callers.
        suffix = _mode_suffix(mode)
        filename = f"{src.stem}{suffix}{ext}"
        if job_id:
            stable_path = jobs_mod.output_path(job_id, filename)
        else:
            stable_out = settings.temp_dir / "out"
            stable_out.mkdir(parents=True, exist_ok=True)
            stable_path = stable_out / filename
        stable_path.unlink(missing_ok=True)
        _link_or_copy(final, stable_path)
        _p(1.0, "Done")
        return stable_path


def _link_or_copy(src: Path, dst: Path) -> None:
    """Hardlink when `src` and `dst` share a filesystem (zero-copy), else
    fall back to copy2. Saves real time on large translated PDFs since both
    the workspace and job dir live under settings.temp_dir."""
    try:
        import os

        os.link(src, dst)
    except (OSError, NotImplementedError):
        shutil.copy2(src, dst)


def _import_pdf_cls():
    from .translators.pdf_translator import PdfTranslator
    return PdfTranslator


def _sample_text(src: Path, fmt: DocFormat) -> str:
    """Grab a short text sample for language detection."""
    return " ".join(sample_paragraphs(src, fmt, max_paragraphs=30))[:2000]


def sample_paragraphs(
    src: Path, fmt: DocFormat, *, max_paragraphs: int = 6, max_chars: int = 260
) -> list[str]:
    """Return the first N non-empty paragraph-like strings from a document.

    Used by the UI's preview endpoint so the progress/result screens can show
    real text from the uploaded and translated artifacts instead of placeholder
    filler. Each entry is trimmed to `max_chars` so the preview stays compact.
    """
    def _finish(items: list[str]) -> list[str]:
        out: list[str] = []
        for raw in items:
            s = " ".join((raw or "").split())
            if not s:
                continue
            if len(s) > max_chars:
                s = s[: max_chars - 1].rstrip() + "…"
            out.append(s)
            if len(out) >= max_paragraphs:
                break
        return out

    try:
        if fmt is DocFormat.TXT:
            text = src.read_text(encoding="utf-8", errors="replace")
            return _finish([p for p in text.split("\n\n") if p.strip()])
        if fmt is DocFormat.DOCX:
            from docx import Document

            d = Document(str(src))
            return _finish([p.text for p in d.paragraphs if p.text and p.text.strip()])
        if fmt is DocFormat.PPTX:
            from pptx import Presentation

            prs = Presentation(str(src))
            parts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        parts.append(shape.text_frame.text)
                if len(parts) >= max_paragraphs * 2:
                    break
            return _finish(parts)
        if fmt is DocFormat.XLSX:
            from openpyxl import load_workbook

            wb = load_workbook(str(src), read_only=True)
            parts: list[str] = []
            try:
                ws = wb.worksheets[0] if wb.worksheets else None
                if ws is None:
                    return []
                for row in ws.iter_rows(max_row=200, values_only=True):
                    row_vals = [str(v) for v in row if isinstance(v, str) and v.strip()]
                    if row_vals:
                        parts.append(" · ".join(row_vals))
                    if len(parts) >= max_paragraphs * 2:
                        break
            finally:
                wb.close()
            return _finish(parts)
        if fmt is DocFormat.PDF:
            import pymupdf

            d = pymupdf.open(str(src))
            try:
                parts: list[str] = []
                for page in d:
                    for block in page.get_text("blocks"):
                        txt = block[4] if len(block) > 4 else ""
                        if isinstance(txt, str) and txt.strip():
                            parts.append(txt)
                    if len(parts) >= max_paragraphs * 2:
                        break
                return _finish(parts)
            finally:
                d.close()
    except Exception:
        return []
    return []


def _mode_suffix(mode: OutputMode) -> str:
    return {
        OutputMode.ORIGINAL: "_original",
        OutputMode.TRANSLATED: "_translated",
        OutputMode.BOTH_VERTICAL: "_both_v",
        OutputMode.BOTH_HORIZONTAL: "_both_h",
    }[mode]
