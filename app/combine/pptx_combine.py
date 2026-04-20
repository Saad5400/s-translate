from __future__ import annotations

import copy
from pathlib import Path


def combine_vertical(src: Path, trans: Path, out: Path) -> Path:
    """Append translated slides after original slides by copying slide XML.

    Approach: load src, iterate translated slides and add each to src. Uses
    the low-level XML-clone trick that python-pptx users commonly employ.
    """
    from pptx import Presentation

    master = Presentation(str(src))
    trans_prs = Presentation(str(trans))

    # XPath to the sldIdLst element inside presentation.xml.
    for slide in trans_prs.slides:
        _clone_slide(master, slide)

    master.save(str(out))
    return out


def _clone_slide(dest_prs, src_slide) -> None:
    """Clone src_slide into dest_prs. Uses deep XML copy."""
    from copy import deepcopy

    # Pick a matching or fallback slide layout in the dest presentation.
    try:
        layout = dest_prs.slide_layouts[0]
    except IndexError:
        layout = src_slide.slide_layout

    new_slide = dest_prs.slides.add_slide(layout)
    # Remove placeholders added by the layout.
    for shp in list(new_slide.shapes):
        sp = shp._element
        sp.getparent().remove(sp)

    # Copy every shape from source slide.
    for shp in src_slide.shapes:
        new_el = deepcopy(shp._element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
