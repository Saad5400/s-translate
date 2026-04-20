from __future__ import annotations

from pathlib import Path

from lxml import etree

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_A = f"{{{_A_NS}}}"


def apply(out_path: Path) -> Path:
    from pptx import Presentation

    prs = Presentation(str(out_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            _apply_shape(shape)
        if slide.has_notes_slide:
            for shape in slide.notes_slide.shapes:
                _apply_shape(shape)
    prs.save(str(out_path))
    return out_path


def _apply_shape(shape) -> None:
    if getattr(shape, "shape_type", None) == 6:  # group
        for sub in shape.shapes:
            _apply_shape(sub)
        return
    if shape.has_text_frame:
        _apply_text_frame(shape.text_frame)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                _apply_text_frame(cell.text_frame)


def _apply_text_frame(tf) -> None:
    for paragraph in tf.paragraphs:
        p = paragraph._pPr_or_insert() if hasattr(paragraph, "_pPr_or_insert") else None
        pPr = paragraph._p.find(f"{_A}pPr")
        if pPr is None:
            pPr = etree.SubElement(paragraph._p, f"{_A}pPr")
            paragraph._p.insert(0, pPr)
        pPr.set("rtl", "1")
        # Right-align
        pPr.set("algn", "r")
