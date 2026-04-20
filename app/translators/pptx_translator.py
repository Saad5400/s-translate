from __future__ import annotations

import re
from pathlib import Path
from typing import Iterable

from ..schemas import DocFormat, Segment
from .base import Translator

_SENTINEL_RE = re.compile(r"⟦(\d+)⟧(.*?)⟦/\1⟧", re.DOTALL)


def _iter_text_frames(shape) -> Iterable:
    """Yield text_frames contained within a shape, recursing into groups/tables."""
    if shape.shape_type == 6:  # GROUP
        for sub in shape.shapes:
            yield from _iter_text_frames(sub)
        return
    if shape.has_text_frame:
        yield shape.text_frame
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                yield cell.text_frame


def _paragraph_joined(paragraph) -> tuple[str, list]:
    runs = list(paragraph.runs)
    if len(runs) == 0:
        return "", runs
    if len(runs) == 1:
        return runs[0].text or "", runs
    parts = [f"⟦{i}⟧{r.text or ''}⟦/{i}⟧" for i, r in enumerate(runs, start=1)]
    return "".join(parts), runs


def _split_translated_into_runs(translated: str, n_runs: int) -> list[str]:
    matches = _SENTINEL_RE.findall(translated)
    if not matches:
        return [translated] + [""] * (n_runs - 1)
    by_idx: dict[int, str] = {}
    for num, text in matches:
        try:
            by_idx[int(num)] = text
        except ValueError:
            continue
    return [by_idx.get(i, "") for i in range(1, n_runs + 1)]


class PptxTranslator(Translator):
    fmt = DocFormat.PPTX

    def extract(self, src_path: Path) -> list[Segment]:
        from pptx import Presentation

        prs = Presentation(str(src_path))
        segments: list[Segment] = []
        idx = 0
        for slide_i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                for tf in _iter_text_frames(shape):
                    for p_i, paragraph in enumerate(tf.paragraphs):
                        joined, runs = _paragraph_joined(paragraph)
                        if not joined.strip():
                            continue
                        segments.append(
                            Segment(
                                id=f"s{idx}",
                                text=joined,
                                meta={"para_path": (slide_i, id(tf), p_i), "n_runs": len(runs)},
                            )
                        )
                        idx += 1
            # Notes slide
            if slide.has_notes_slide:
                tf = slide.notes_slide.notes_text_frame
                for p_i, paragraph in enumerate(tf.paragraphs):
                    joined, runs = _paragraph_joined(paragraph)
                    if not joined.strip():
                        continue
                    segments.append(
                        Segment(
                            id=f"s{idx}",
                            text=joined,
                            meta={"para_path": ("notes", slide_i, p_i), "n_runs": len(runs)},
                        )
                    )
                    idx += 1
        return segments

    def reinsert(self, src_path: Path, segments: list[Segment], out_path: Path) -> Path:
        from pptx import Presentation

        prs = Presentation(str(src_path))
        # Re-walk and apply translations in the same order.
        translations: list[str] = [s.translated or s.text for s in segments]
        it = iter(translations)

        def apply(paragraph, runs: list) -> None:
            try:
                tr = next(it)
            except StopIteration:
                return
            n = len(runs)
            if n == 0:
                return
            if n == 1:
                runs[0].text = _SENTINEL_RE.sub(lambda m: m.group(2), tr)
            else:
                parts = _split_translated_into_runs(tr, n)
                non_empty = sum(1 for p in parts if p)
                if non_empty <= 1:
                    whole = _SENTINEL_RE.sub(lambda m: m.group(2), tr)
                    runs[0].text = whole
                    for r in runs[1:]:
                        r.text = ""
                else:
                    for r, part in zip(runs, parts, strict=True):
                        r.text = part

        for slide in prs.slides:
            for shape in slide.shapes:
                for tf in _iter_text_frames(shape):
                    for paragraph in tf.paragraphs:
                        joined, runs = _paragraph_joined(paragraph)
                        if not joined.strip():
                            continue
                        apply(paragraph, runs)
            if slide.has_notes_slide:
                tf = slide.notes_slide.notes_text_frame
                for paragraph in tf.paragraphs:
                    joined, runs = _paragraph_joined(paragraph)
                    if not joined.strip():
                        continue
                    apply(paragraph, runs)

        prs.save(str(out_path))
        return out_path
