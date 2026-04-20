"""Programmatically create sample docs in each format for integration tests.

Run directly:
    .venv/bin/python -m tests.make_fixtures
"""
from __future__ import annotations

from pathlib import Path

FIXTURES = Path(__file__).parent / "fixtures"
FIXTURES.mkdir(parents=True, exist_ok=True)

TITLE = "Quarterly Report"
PARAGRAPHS = [
    "This document summarises quarterly performance.",
    "Revenue grew by 12% compared to the prior quarter.",
    "Operating expenses remained flat thanks to efficiency improvements.",
    "We anticipate continued growth in the upcoming quarter.",
]
BULLETS = [
    "Increase marketing spend in key regions.",
    "Launch the new product line in May.",
    "Expand the engineering team by 20%.",
]
TABLE = [
    ["Metric", "Q1", "Q2"],
    ["Revenue", "100", "112"],
    ["Expenses", "80", "81"],
    ["Profit", "20", "31"],
]


def make_txt() -> Path:
    p = FIXTURES / "sample.txt"
    content = (
        TITLE
        + "\n\n"
        + "\n\n".join(PARAGRAPHS)
        + "\n\n"
        + "Key priorities:\n"
        + "\n".join(f"- {b}" for b in BULLETS)
    )
    p.write_text(content, encoding="utf-8")
    return p


def make_docx() -> Path:
    from docx import Document
    from docx.shared import Pt, RGBColor

    doc = Document()
    h = doc.add_heading(TITLE, level=1)
    for para in PARAGRAPHS:
        p = doc.add_paragraph(para)
    # Add a paragraph with mixed run formatting (bold + italic runs).
    p = doc.add_paragraph()
    r1 = p.add_run("Important: ")
    r1.bold = True
    r2 = p.add_run("please review the ")
    r3 = p.add_run("quarterly figures")
    r3.italic = True
    p.add_run(" before the meeting.")

    doc.add_heading("Key priorities", level=2)
    for b in BULLETS:
        doc.add_paragraph(b, style="List Bullet")

    # Table
    doc.add_heading("Financials", level=2)
    table = doc.add_table(rows=len(TABLE), cols=len(TABLE[0]))
    table.style = "Light Grid Accent 1"
    for i, row in enumerate(TABLE):
        for j, cell in enumerate(row):
            table.rows[i].cells[j].text = cell

    p = FIXTURES / "sample.docx"
    doc.save(str(p))
    return p


def make_pptx() -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
    slide.shapes.title.text = TITLE
    slide.placeholders[1].text = "Prepared by Finance Team"

    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + content
    slide.shapes.title.text = "Summary"
    tf = slide.placeholders[1].text_frame
    tf.text = PARAGRAPHS[0]
    for para in PARAGRAPHS[1:]:
        p = tf.add_paragraph()
        p.text = para

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key priorities"
    tf = slide.placeholders[1].text_frame
    tf.text = BULLETS[0]
    for b in BULLETS[1:]:
        p = tf.add_paragraph()
        p.text = b

    p = FIXTURES / "sample.pptx"
    prs.save(str(p))
    return p


def make_xlsx() -> Path:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill

    wb = Workbook()
    ws = wb.active
    ws.title = "Report"

    ws["A1"] = TITLE
    ws["A1"].font = Font(bold=True, size=14)

    for i, row in enumerate(TABLE, start=3):
        for j, cell in enumerate(row, start=1):
            c = ws.cell(row=i, column=j, value=cell)
            if i == 3:  # header row
                c.font = Font(bold=True)
                c.fill = PatternFill("solid", fgColor="DDDDDD")
                c.alignment = Alignment(horizontal="center")

    # Priorities sheet
    ws2 = wb.create_sheet("Priorities")
    ws2["A1"] = "Key priorities"
    ws2["A1"].font = Font(bold=True)
    for i, b in enumerate(BULLETS, start=2):
        ws2.cell(row=i, column=1, value=b)

    p = FIXTURES / "sample.xlsx"
    wb.save(p)
    return p


def make_pdf() -> Path:
    """Build a styled PDF via PyMuPDF's HTML box API."""
    import pymupdf

    doc = pymupdf.open()
    page = doc.new_page()
    html_content = (
        f"<h1 style='font-family: Helvetica;'>{TITLE}</h1>"
        + "".join(f"<p style='font-family: Helvetica; font-size: 11pt;'>{p}</p>" for p in PARAGRAPHS)
        + "<h2 style='font-family: Helvetica;'>Key priorities</h2>"
        + "<ul style='font-family: Helvetica; font-size: 11pt;'>"
        + "".join(f"<li>{b}</li>" for b in BULLETS)
        + "</ul>"
    )
    page.insert_htmlbox(pymupdf.Rect(40, 40, 555, 750), html_content)

    # Second page with a table rendered via HTML.
    page2 = doc.new_page()
    rows_html = "".join(
        "<tr>" + "".join(f"<td style='padding:4px;border:1px solid #888'>{c}</td>" for c in row) + "</tr>"
        for row in TABLE
    )
    page2.insert_htmlbox(
        pymupdf.Rect(40, 40, 555, 750),
        f"<h2>Financials</h2><table style='border-collapse:collapse;font-family:Helvetica'>{rows_html}</table>",
    )

    p = FIXTURES / "sample.pdf"
    doc.save(str(p))
    doc.close()
    return p


def make_all() -> dict[str, Path]:
    return {
        "txt": make_txt(),
        "docx": make_docx(),
        "pptx": make_pptx(),
        "xlsx": make_xlsx(),
        "pdf": make_pdf(),
    }


if __name__ == "__main__":
    paths = make_all()
    for fmt, p in paths.items():
        print(f"{fmt}: {p} ({p.stat().st_size} bytes)")
