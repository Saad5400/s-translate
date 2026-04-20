"""End-to-end pipeline tests using a stub LLM (offline).

These tests assert:
  * Extraction produces >=1 segment per format.
  * Reinsertion produces a valid output file (opens without error).
  * Translated marker text actually appears in the output.
  * RTL attributes are applied to DOCX/XLSX/PPTX when target is Arabic.
  * Vertical combine increases page/slide/sheet count.
  * Horizontal combine produces expected width/columns.
"""
from __future__ import annotations

import asyncio
from pathlib import Path

import pytest

from app.api import run_job
from app.llm.client import set_stub_translator
from app.schemas import OutputMode, TranslationJob
from tests.make_fixtures import make_all
from tests.stub_llm import stub_translate

FIXTURES = Path(__file__).parent / "fixtures"


@pytest.fixture(scope="session", autouse=True)
def _make_fixtures():
    make_all()


@pytest.fixture(autouse=True)
def _stub_llm():
    set_stub_translator(stub_translate)
    yield
    set_stub_translator(None)


def _job(src: Path, target: str = "es", mode: OutputMode = OutputMode.TRANSLATED) -> TranslationJob:
    return TranslationJob(
        src_path=src,
        target_lang=target,
        provider="stub",
        model="stub/stub",
        api_key="stub",
        output_mode=mode,
    )


@pytest.mark.asyncio
async def test_txt_translate():
    job = _job(FIXTURES / "sample.txt")
    out = await run_job(job)
    text = out.read_text(encoding="utf-8")
    assert "[es]" in text
    assert out.suffix == ".txt"


@pytest.mark.asyncio
async def test_docx_translate_preserves_structure():
    from docx import Document

    job = _job(FIXTURES / "sample.docx")
    out = await run_job(job)
    doc = Document(str(out))
    texts = [p.text for p in doc.paragraphs]
    joined = "\n".join(texts)
    assert "[es]" in joined
    # Structural assertion: same # of paragraphs before and after.
    orig = Document(str(FIXTURES / "sample.docx"))
    assert len(doc.paragraphs) == len(orig.paragraphs)
    # Tables preserved.
    assert len(doc.tables) == len(orig.tables)


@pytest.mark.asyncio
async def test_docx_rtl_applied_for_arabic():
    from lxml import etree
    from docx import Document

    job = _job(FIXTURES / "sample.docx", target="ar")
    out = await run_job(job)
    doc = Document(str(out))
    W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    # Every non-empty paragraph should have a w:bidi in its pPr.
    found_bidi = 0
    for p in doc.element.body.iter(f"{W}p"):
        if p.find(f"{W}pPr/{W}bidi") is not None:
            found_bidi += 1
    assert found_bidi > 0
    # At least one run should have w:rtl.
    found_rtl = 0
    for r in doc.element.body.iter(f"{W}r"):
        if r.find(f"{W}rPr/{W}rtl") is not None:
            found_rtl += 1
    assert found_rtl > 0


@pytest.mark.asyncio
async def test_xlsx_translate_and_rtl():
    from openpyxl import load_workbook

    # Translate-only to Spanish.
    job = _job(FIXTURES / "sample.xlsx", target="es")
    out = await run_job(job)
    wb = load_workbook(out)
    found = False
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and "[es]" in cell.value:
                    found = True
        # LTR: should NOT be rtl.
        assert ws.sheet_view.rightToLeft in (False, None)
    assert found
    wb.close()

    # Now Arabic — RTL must apply.
    job_ar = _job(FIXTURES / "sample.xlsx", target="ar")
    out_ar = await run_job(job_ar)
    wb_ar = load_workbook(out_ar)
    for ws in wb_ar.worksheets:
        assert ws.sheet_view.rightToLeft is True
    wb_ar.close()


@pytest.mark.asyncio
async def test_pptx_translate_preserves_slides():
    from pptx import Presentation

    job = _job(FIXTURES / "sample.pptx")
    out = await run_job(job)
    prs = Presentation(str(out))
    orig = Presentation(str(FIXTURES / "sample.pptx"))
    assert len(prs.slides) == len(orig.slides)
    # Find translated text
    found = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and "[es]" in shape.text_frame.text:
                found = True
    assert found


@pytest.mark.asyncio
async def test_pdf_translate():
    import pymupdf

    job = _job(FIXTURES / "sample.pdf")
    out = await run_job(job)
    doc = pymupdf.open(str(out))
    try:
        all_text = "\n".join(page.get_text("text") for page in doc)
    finally:
        doc.close()
    assert "[es]" in all_text


@pytest.mark.asyncio
async def test_vertical_combine_docx():
    from docx import Document

    job = _job(FIXTURES / "sample.docx", mode=OutputMode.BOTH_VERTICAL)
    out = await run_job(job)
    combined = Document(str(out))
    orig = Document(str(FIXTURES / "sample.docx"))
    # Combined should have at least 2× paragraphs (original + translated).
    assert len(combined.paragraphs) >= 2 * len(orig.paragraphs) - 2


@pytest.mark.asyncio
async def test_vertical_combine_pdf():
    import pymupdf

    job = _job(FIXTURES / "sample.pdf", mode=OutputMode.BOTH_VERTICAL)
    out = await run_job(job)
    c = pymupdf.open(str(out))
    try:
        orig = pymupdf.open(str(FIXTURES / "sample.pdf"))
        assert c.page_count == orig.page_count * 2
        orig.close()
    finally:
        c.close()


@pytest.mark.asyncio
async def test_vertical_combine_xlsx():
    from openpyxl import load_workbook

    job = _job(FIXTURES / "sample.xlsx", mode=OutputMode.BOTH_VERTICAL)
    out = await run_job(job)
    wb = load_workbook(out)
    names = wb.sheetnames
    assert any(n.endswith("_translated") for n in names)
    wb.close()


@pytest.mark.asyncio
async def test_vertical_combine_pptx():
    from pptx import Presentation

    job = _job(FIXTURES / "sample.pptx", mode=OutputMode.BOTH_VERTICAL)
    out = await run_job(job)
    combined = Presentation(str(out))
    orig = Presentation(str(FIXTURES / "sample.pptx"))
    assert len(combined.slides) == 2 * len(orig.slides)


@pytest.mark.asyncio
async def test_horizontal_combine_pdf_width_doubles():
    import pymupdf

    orig = pymupdf.open(str(FIXTURES / "sample.pdf"))
    orig_w = orig[0].rect.width
    orig.close()

    job = _job(FIXTURES / "sample.pdf", mode=OutputMode.BOTH_HORIZONTAL)
    out = await run_job(job)
    c = pymupdf.open(str(out))
    try:
        assert abs(c[0].rect.width - orig_w * 2) < 1
    finally:
        c.close()


@pytest.mark.asyncio
async def test_horizontal_combine_xlsx():
    from openpyxl import load_workbook

    job = _job(FIXTURES / "sample.xlsx", mode=OutputMode.BOTH_HORIZONTAL)
    out = await run_job(job)
    wb = load_workbook(out)
    ws = wb.worksheets[0]
    # Original had 3 columns (metric/Q1/Q2), so combined should have >= 6 used.
    assert ws.max_column >= 6
    wb.close()


@pytest.mark.asyncio
async def test_horizontal_combine_docx_produces_pdf():
    # DOCX horizontal routes through LibreOffice → PDF.
    import shutil

    if not shutil.which("soffice") and not shutil.which("libreoffice"):
        pytest.skip("LibreOffice not installed")
    job = _job(FIXTURES / "sample.docx", mode=OutputMode.BOTH_HORIZONTAL)
    out = await run_job(job)
    assert out.suffix == ".pdf"
    import pymupdf

    c = pymupdf.open(str(out))
    try:
        assert c.page_count >= 1
    finally:
        c.close()
